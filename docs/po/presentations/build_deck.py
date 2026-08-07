#!/usr/bin/env python3
"""Generate the LifeLink KH Week-2 pitch deck (version 1).

Run from anywhere:
    python3 docs/po/presentations/build_deck.py

Output: docs/po/presentations/LifeLinkKH-v1.pptx

The milestone slide is parsed from CLAUDE.md section 4 at generation time rather than
hardcoded. Milestone assignments already changed once (DEC-001/002/003 in docs/decisions.md);
a hardcoded copy here would be another place that table lives and the first to go stale. Re-run this script after any milestone change.
"""

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[3]
OUT = Path(__file__).resolve().parent / "LifeLinkKH-v1.pptx"

# --- Design tokens -----------------------------------------------------------
# Dark text on light background: this is shown on a washed-out classroom projector,
# where light-on-dark loses all contrast.
INK = RGBColor(0x1A, 0x1A, 0x2E)      # near-black body text
ACCENT = RGBColor(0xC0, 0x21, 0x2E)   # blood red, used sparingly
TEAL = RGBColor(0x0F, 0x6B, 0x5F)     # secondary accent
MUTED = RGBColor(0x5A, 0x5A, 0x6E)    # captions, footers
BG = RGBColor(0xFA, 0xFA, 0xF7)       # off-white slide background

LATIN_FONT = "Helvetica Neue"
# Noto Sans Khmer, not Khmer Sangam MN: Sangam is macOS-only and renders as boxes on a
# Windows classroom machine. Noto ships with Google Slides and most Linux distros.
KHMER_FONT = "Noto Sans Khmer"

TITLE_PT = 34
KHMER_TITLE_PT = 30
BODY_PT = 19
SMALL_PT = 15

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# --- Helpers -----------------------------------------------------------------
def read(rel):
    return (REPO / rel).read_text(encoding="utf-8")


def parse_milestones():
    """Pull the M1-M7 table out of CLAUDE.md section 4.

    Returns [(id, weeks, deliverable), ...]. Raises if the table cannot be found, so a
    silent empty slide is impossible.
    """
    text = read("CLAUDE.md")
    rows = re.findall(r"^\|\s*(M[1-7])\s*\|\s*([^|]+?)\s*\|\s*(.+?)\s*\|\s*$",
                      text, re.MULTILINE)
    if len(rows) != 7:
        sys.exit(f"ERROR: expected 7 milestone rows in CLAUDE.md section 4, found {len(rows)}. "
                 "Fix the table or this script's regex before generating the deck.")
    # Strip markdown emphasis so **bold** does not leak onto the slide.
    return [(m, w, re.sub(r"\*\*|`", "", d)) for m, w, d in rows]


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide


def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def style(run, size, color=INK, bold=False, khmer=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = KHMER_FONT if khmer else LATIN_FONT


def heading(slide, khmer_text, english_text, kicker=None):
    """Khmer title with the English underneath, per the deck's bilingual rule."""
    tf = textbox(slide, Inches(0.8), Inches(0.45), Inches(11.7), Inches(1.7))

    if kicker:
        p = tf.paragraphs[0]
        style(p.add_run(), SMALL_PT)
        p.runs[0].text = kicker.upper()
        p.runs[0].font.color.rgb = TEAL
        p.runs[0].font.bold = True
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]

    run = p.add_run()
    run.text = khmer_text
    style(run, KHMER_TITLE_PT, INK, bold=True, khmer=True)

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = english_text
    style(run2, TITLE_PT, ACCENT, bold=True)

    # Rule under the heading
    line = slide.shapes.add_shape(1, Inches(0.85), Inches(2.28), Inches(1.4), Emu(34925))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    line.shadow.inherit = False
    return tf


def bullets(slide, items, top=Inches(2.7), left=Inches(0.9), width=Inches(11.5),
            size=BODY_PT):
    """Max 6 items, max 12 words each — detail belongs in the speaker notes.

    Punctuation-only tokens (separators like the interpunct and em dash) are not words and
    do not count against the cap.
    """
    assert len(items) <= 6, f"{len(items)} bullets: cap is 6"
    for text in items:
        words = [t for t in text.split() if any(c.isalnum() for c in t)]
        assert len(words) <= 12, f"bullet over 12 words ({len(words)}): {text!r}"

    tf = textbox(slide, left, top, width, Inches(4.0))
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        dot = p.add_run()
        dot.text = "▪  "
        style(dot, size, ACCENT, bold=True)
        run = p.add_run()
        run.text = text
        style(run, size)
    return tf


def footer(slide, text):
    tf = textbox(slide, Inches(0.85), Inches(6.85), Inches(11.6), Inches(0.4))
    run = tf.paragraphs[0].add_run()
    run.text = text
    style(run, 12, MUTED)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


FOOT = "LifeLink KH · Group 2 · Track B · Week 2 of 16"


# --- Slides ------------------------------------------------------------------
def slide_01_title(prs):
    s = blank(prs)
    tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(3.2))

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "ជីវិត"
    style(r, 60, ACCENT, bold=True, khmer=True)

    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "LifeLink KH"
    style(r, 46, INK, bold=True)

    p = tf.add_paragraph()
    p.space_before = Pt(16)
    r = p.add_run()
    r.text = "Blood donor matching for Cambodia"
    style(r, 24, MUTED)

    tf2 = textbox(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.2))
    p = tf2.paragraphs[0]
    r = p.add_run()
    r.text = "Nem Sothea · Suon Pisey · Sourn SAVOURN · Moeun Nithvaraman · Oun Sreynich"
    style(r, 16, INK)
    p2 = tf2.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Group 2 · Cross-Platform Mobile App Development · Track B · Week 2 of 16"
    style(r2, 14, MUTED)

    notes(s, """
Open with the problem, not the product. One line: "Right now, when a Cambodian hospital
needs blood urgently, someone writes a Facebook post and hopes the right person scrolls
past." Then move on — slide 2 does the work.

Do not read the team list aloud; it is on the slide and slide 12 covers roles.
""")
    return s


def slide_02_problem(prs):
    s = blank(prs)
    heading(s, "បញ្ហា", "The Problem", kicker="Why")
    bullets(s, [
        "Hospitals and the National Blood Transfusion Center rely on Facebook posts",
        "A post reaches whoever scrolls past — not compatible donors nearby",
        "No way to filter by blood type, distance, or donor eligibility",
        "No way to know if anyone is coming to help",
        "Blood emergencies are measured in minutes; this process is not",
    ])
    footer(s, FOOT)
    notes(s, """
This is the slide that earns attention. Stay concrete and do not inflate it.

Say: the current process is not a bad app, it is no app. It is a social media post and
hope. Reach without targeting is not reach — most people who see the post have the wrong
blood type or are in the wrong city, and the people who could actually help never see it.

If asked for national statistics: we have not sourced verified figures yet, and we would
rather present none than present numbers we cannot cite. The problem statement comes from
how hospitals currently operate, which is observable without a statistic.
""")
    return s


def slide_03_users(prs):
    s = blank(prs)
    heading(s, "អ្នកប្រើប្រាស់គោលដៅ", "Target Users", kicker="Who")
    bullets(s, [
        "Donor — a voluntary blood donor, on the mobile app",
        "Requester — a patient or family member needing blood, mobile",
        "Hospital staff — verifies requests, confirms donations, on web",
        "Administrator — manages users and hospitals, on web",
        "Pilot scope: Phnom Penh first, other provinces later",
    ])
    footer(s, FOOT)
    notes(s, """
The assignment asks for a target user. Ours is the Donor — everything else exists to serve
that person's decision to help.

Emphasise the split: donors and families are on the phone because emergencies happen away
from a desk. Hospital staff are on the web because they work at a desk with a browser. That
split is why the project has two clients, which slide 6 explains.

The four roles come from the PRD, section 3. Do not invent a fifth on the spot.
""")
    return s


def slide_04_features(prs):
    s = blank(prs)
    heading(s, "មុខងារសំខាន់បី", "Three Core Features", kicker="What")
    bullets(s, [
        "1 · Donor register — blood type, location, automatic 56-day eligibility check",
        "2 · Urgent request broadcast — push to matching donors by type and distance",
        "3 · History + reminder — tracks cooldown, alerts donor when eligible again",
        "Matching uses ABO/Rh compatibility, not exact blood type only",
        "Khmer and English on both clients, Khmer as the default",
    ])
    footer(s, FOOT)
    notes(s, """
Exactly three features, as the assignment requires. Say the numbers out loud — "one, two,
three" — so the grader can tick them off.

The fourth bullet is worth pausing on: O-negative can help almost anyone, AB-positive can
receive from anyone. Matching on exact type alone would throw away most of the available
supply. That detail signals we understood the domain rather than guessing at it.

If asked what is NOT in version 1: in-app chat, rewards or gamification, hospital blood
inventory, iOS, and automated medical verification. All are recorded as out of scope in the
PRD, section 2.2 — deliberately excluded, not forgotten.
""")
    return s


def slide_05_why_mobile(prs):
    s = blank(prs)
    heading(s, "ហេតុអ្វីត្រូវជាកម្មវិធីទូរស័ព្ទ", "Why Mobile, Not a Website", kicker="Why")

    tf = textbox(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(2.0))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ("A blood emergency needs an instant, location-aware push alert on a donor's "
              "phone — a website can neither push a time-critical notification nor read "
              "real-time GPS.")
    style(r, 27, INK, bold=True)

    tf2 = textbox(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.6))
    for i, text in enumerate([
        "Push notification — reaches a donor who is not looking at their phone",
        "Real-time GPS — ranks donors by actual distance to the hospital",
        "Installed app — no URL to find while a relative is in a corridor",
    ]):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(8)
        dot = p.add_run()
        dot.text = "▪  "
        style(dot, SMALL_PT, TEAL, bold=True)
        r = p.add_run()
        r.text = text
        style(r, SMALL_PT, MUTED)

    footer(s, FOOT)
    notes(s, """
This is a graded requirement — one sentence, defended. Deliver the sentence in bold on the
slide verbatim, then stop. Do not elaborate unless asked.

LIKELY QUESTION: "Couldn't a progressive web app do push notifications?"
ANSWER: Partly, and that is the honest answer. Web push exists, but it is unreliable on the
devices our donors actually use, it cannot read GPS in the background, and it cannot be
distributed through the Play Store — which is this course's M7 requirement. We also need the
app installed and present on the home screen, because a donor who has to remember a URL
during an emergency will not use it.

LIKELY QUESTION: "Then why build a web portal at all?"
ANSWER: Different user, different context. Hospital staff work at a desk and need to
confirm donations; that is a browser job. The mobile app is for the people who are moving.
""")
    return s


def slide_06_architecture(prs):
    s = blank(prs)
    heading(s, "ស្ថាបត្យកម្មប្រព័ន្ធ", "Architecture", kicker="How")

    diagram = (
        "                Spring Boot API  ──▶  PostgreSQL (Flyway)\n"
        "                    ▲          ▲\n"
        "       REST/HTTPS   │          │   REST/HTTPS\n"
        "                    │          │\n"
        "       Flutter app ─┘          └─ Next.js web portal\n"
        "    (donors / patients)          (hospitals / admin)\n"
        "         → Play Store"
    )
    tf = textbox(s, Inches(0.9), Inches(2.65), Inches(6.6), Inches(3.0))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = diagram
    r.font.size = Pt(13)
    r.font.name = "Menlo"
    r.font.color.rgb = INK

    tf2 = textbox(s, Inches(7.8), Inches(2.65), Inches(4.7), Inches(3.6))
    rows = [
        ("Backend", "Spring Boot · PostgreSQL · Flyway"),
        ("Web", "Next.js · TypeScript · Tailwind"),
        ("Mobile", "Flutter → native Android"),
        ("Push", "Firebase Cloud Messaging"),
        ("Location", "Google Maps · geolocator"),
        ("Local dev + CI", "Docker Compose · GitHub Actions"),
    ]
    for i, (k, v) in enumerate(rows):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(11)
        kr = p.add_run()
        kr.text = f"{k}   "
        style(kr, SMALL_PT, TEAL, bold=True)
        vr = p.add_run()
        vr.text = v
        style(vr, SMALL_PT)

    footer(s, FOOT)
    notes(s, """
One API, two clients. Say that first; the diagram then reads itself.

The key decision here is Flutter rather than a hybrid wrapper. Flutter builds a native
Android binary, which is what the Play Store requirement at M7 actually needs. We considered
wrapping the Next.js app in Capacitor and rejected it — that is recorded in ADR 0001.

Flyway matters more than it looks: the database schema is versioned migrations, not
hand-edited tables, so every environment gets the same schema in the same order.

If asked why not one codebase for everything: the two clients serve different users in
different contexts, and they map cleanly onto team roles.
""")
    return s


def slide_07_scope(prs):
    s = blank(prs)
    heading(s, "ការគ្រប់គ្រងវិសាលភាព", "How We Manage Scope", kicker="How")
    bullets(s, [
        "19 features registered — 8 in the build, 8 deliberately deferred",
        "Cut against what the course grades, not what interests us",
        "Definition of Done: spec signed off, merged, QA verified, no open bugs",
        "Every architecture choice recorded as a decision record",
        "Deferred features keep their documents — that is our future work",
    ])
    footer(s, FOOT)
    notes(s, """
This slide answers a question the audience has not asked yet: is this a plan or a wish?

Two things to say here, in this order.

First, we read our own PRD adversarially and found seven promises it made in prose but never
turned into requirements — for example, it said users may request deletion of their data, and
no feature provided it. That is how the registry got to 19.

Second, we then cut it to 8. Nineteen features across thirteen part-time weeks does not fit,
and a plan that does not fit produces eight half-finished features instead of eight working
ones. So we cut against the five things the course grades — authentication, push, GPS, a
relational database, a Play Store release — not against what we personally found interesting.

LIKELY QUESTION: "Why isn't feature X in your app?"
ANSWER: It probably is in docs/scope.md, with the reason. Name one deliberately: account and
data deletion is deferred, and we flag it as a privacy obligation rather than a feature — it
is only acceptable because the pilot runs on our own test accounts, and it comes back first
if a real donor ever uses this.

If asked how we track work: features have stable IDs, so a bug, a test case and a commit can
all point at the same identifier. QA signs off against acceptance criteria, not against a
demo.

Do not spend more than a minute here. It is a credibility slide, not the substance.
""")
    return s


def slide_08_milestones(prs, milestones):
    s = blank(prs)
    heading(s, "ដំណាក់កាលការងារ", "Milestones · M1 → M7", kicker="When")

    tf = textbox(s, Inches(0.9), Inches(2.55), Inches(11.6), Inches(4.1))
    for i, (mid, weeks, deliverable) in enumerate(milestones):
        text = deliverable if len(deliverable) <= 96 else deliverable[:93].rstrip() + "…"
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        final = mid == "M7"

        mr = p.add_run()
        mr.text = f"{mid}  "
        style(mr, SMALL_PT + 2, ACCENT if final else TEAL, bold=True)

        wr = p.add_run()
        wr.text = f"{weeks.strip():8s}  "
        style(wr, SMALL_PT, MUTED)

        dr = p.add_run()
        dr.text = text
        style(dr, SMALL_PT, INK, bold=final)

    footer(s, "Parsed from CLAUDE.md section 4 — the single source for milestone dates")
    notes(s, """
Do not read all seven aloud. Name three: M2 is the stack running under Docker, M4 is the
core matching feature end to end, M7 is published to Play Store internal testing by Week 15.

M7 is the assignment's hard requirement, so say the words "internal testing" explicitly.

This table is generated from our planning document, not retyped — which is why it cannot
disagree with the plan. We changed three milestone rows this week when we found that
matching was scheduled before the eligibility check it depends on. Slide 11 covers that.
""")
    return s


def slide_09_status(prs):
    s = blank(prs)
    heading(s, "ស្ថានភាពបច្ចុប្បន្ន", "Where We Are Today", kicker="When")
    bullets(s, [
        "M1 complete — PRD, feature registry, schema design, API contract skeleton",
        "M2 specified in four layer specs, each with a pass/fail checklist",
        "Not yet built — the three code directories are still empty",
        "Blocking M2: one API endpoint must reach the shared contract",
        "Next: resolve three product decisions that gate M4 matching",
    ])
    footer(s, FOOT)
    notes(s, """
Be straight about this. We have documentation and specifications, not running code. Week 2
of 16 is the right time to have exactly that.

What the M2 specs contain: the Flutter project structure and state-management choice, the
Spring Boot package layout and the initial database migration covering all six entities,
the Next.js routing and localization approach, and the Docker Compose service definitions.
Each has a checklist a reviewer can pass or fail without asking us anything.

LIKELY QUESTION: "Can you show us the app?"
ANSWER: No, and we would rather say so than show a mockup that implies working software.
What we can show is the specification the build will be checked against.
""")
    return s


def slide_10_metrics(prs):
    s = blank(prs)
    heading(s, "សូចនាករជោគជ័យ", "Success Metrics", kicker="Measure")
    bullets(s, [
        "≥ 200 registered donors in Phnom Penh in the first pilot month",
        "≥ 70% of urgent requests get an acceptance within 60 minutes",
        "Median time from request to first acceptance under 30 minutes",
        "≥ 50 completed, hospital-verified donations during the pilot semester",
        "Push notification delivery rate ≥ 95%",
    ])
    footer(s, FOOT)
    notes(s, """
Five numbers, all from the PRD. These are targets we set, not results we have.

The important one is the third: median time to first acceptance. It is the number that says
whether we actually beat a Facebook post, and it is the reason the app exists.

Note that hospital-verified matters in the fourth metric. A self-reported donation is not
evidence; a donation confirmed by hospital staff is.

LIKELY QUESTION: "How will you measure these?"
ANSWER: With SQL COUNT queries against the pilot database at demo time. We originally planned
event-capture instrumentation on every milestone and cut it — for a pilot this size the
queries give the same five numbers for none of the build cost. Be honest about the limit: if
pilot data is thin, these stay targets rather than results.
""")
    return s


def slide_11_risks(prs):
    s = blank(prs)
    heading(s, "ហានិភ័យ និងការសម្រេចចិត្ត", "Risks + Open Decisions", kicker="Honest")
    bullets(s, [
        "Two decisions still open: notified-donor count, deploy runbook",
        "M4 still heaviest — now three weeks, cut order agreed",
        "Low donor density — zero-match fallback deferred, so unmitigated",
        "Auth switched to Google Sign-In — phone numbers now unverified",
        "Data deletion deferred — must ship before any real donor",
    ])
    footer(s, FOOT)
    notes(s, """
Most Week-2 decks hide what they do not know. This slide is deliberate — it is the one that
shows the plan is real.

On location precision: storing exact coordinates gives accurate distance ranking but holds
sensitive personal data; a district centroid is safer but less accurate. It is simultaneously
a privacy decision and a product decision, and it changes the database schema, so we are
deciding it before we write the migration rather than after.

On M4: we found that matching was scheduled before the eligibility check it depends on, and
that push notification was scheduled after the feature that triggers it. We moved both
earlier rather than let a milestone close with its own acceptance criteria unmet. That made
M4 heavy, so we agreed in advance which two features get cut first if it slips.

LIKELY QUESTION: "Isn't it a problem that so much is undecided?"
ANSWER: These are written down with the milestone each one blocks and the date it must be
resolved. The alternative is not fewer unknowns — it is the same unknowns, undocumented.
""")
    return s


def slide_12_team(prs):
    s = blank(prs)
    heading(s, "ក្រុមការងារ", "Team + One Open Question", kicker="Who builds it")

    tf = textbox(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.9))
    team = [
        ("Nem Sothea", "Tech Lead · Flutter · also PO, Security · Docker, CI, release"),
        ("Suon Pisey", "Backend · PostgreSQL · Flyway · matching logic"),
        ("Sourn SAVOURN", "Frontend · Next.js portal · API client · i18n"),
        ("Moeun Nithvaraman", "PO · PRD · briefs · wireframes · feature registry"),
        ("Oun Sreynich", "QA · test plan · milestone acceptance · bug tracking"),
    ]
    for i, (name, role) in enumerate(team):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        nr = p.add_run()
        nr.text = f"{name}   "
        style(nr, BODY_PT, INK, bold=True)
        rr = p.add_run()
        rr.text = role
        style(rr, SMALL_PT, MUTED)

    box = s.shapes.add_shape(1, Inches(0.85), Inches(5.75), Inches(11.6), Inches(0.85))
    box.fill.solid()
    box.fill.fore_color.rgb = TEAL
    box.line.fill.background()
    box.shadow.inherit = False
    btf = box.text_frame
    btf.word_wrap = True
    p = btf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = ("Open question for the instructor: the assignment specifies teams of 3 — "
              "we are 5. How should we proceed?")
    r.font.size = Pt(17)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r.font.name = LATIN_FONT

    footer(s, FOOT)
    notes(s, """
End on the question, not on thanks. Ask it directly and then stop talking.

We are raising it rather than waiting to be asked, because the answer may change team
composition or how the work is graded, and it is cheaper to know in Week 2 than Week 10.

If the answer is "split into teams of 3": our roles already divide along clean boundaries —
each person owns a directory and a document set, so a split is a reassignment, not a rewrite.

If asked who does what day to day: there is no separate DevOps or PM role. Tech Lead absorbed that
work — Docker, CI, deploy, release. Definition-of-Done tracking deliberately stayed with QA, because
with Tech Lead also holding Security and co-PO, QA is the only gate outside one person.
We know the gap: there is no deploy runbook yet, and it has to exist before the M7 release. Product definition is co-held by Moeun and Sothea, so PRD and FR sign-off is
not one person's signature. Tech Lead still holds Security, which concentrates approval there;
QA sign-off is kept independent for exactly that reason.
""")
    return s


# --- Main --------------------------------------------------------------------
def main():
    milestones = parse_milestones()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        ("1", "ជីវិត · LifeLink KH", lambda: slide_01_title(prs)),
        ("2", "បញ្ហា · The Problem", lambda: slide_02_problem(prs)),
        ("3", "អ្នកប្រើប្រាស់គោលដៅ · Target Users", lambda: slide_03_users(prs)),
        ("4", "មុខងារសំខាន់បី · Three Core Features", lambda: slide_04_features(prs)),
        ("5", "ហេតុអ្វីត្រូវជាកម្មវិធីទូរស័ព្ទ · Why Mobile", lambda: slide_05_why_mobile(prs)),
        ("6", "ស្ថាបត្យកម្មប្រព័ន្ធ · Architecture", lambda: slide_06_architecture(prs)),
        ("7", "ការគ្រប់គ្រងវិសាលភាព · How We Manage Scope", lambda: slide_07_scope(prs)),
        ("8", "ដំណាក់កាលការងារ · Milestones", lambda: slide_08_milestones(prs, milestones)),
        ("9", "ស្ថានភាពបច្ចុប្បន្ន · Where We Are Today", lambda: slide_09_status(prs)),
        ("10", "សូចនាករជោគជ័យ · Success Metrics", lambda: slide_10_metrics(prs)),
        ("11", "ហានិភ័យ · Risks + Open Decisions", lambda: slide_11_risks(prs)),
        ("12", "ក្រុមការងារ · Team + One Open Question", lambda: slide_12_team(prs)),
    ]
    for num, label, build in builders:
        build()
        print(f"OK slide {num} - {label}")

    prs.save(OUT)
    print(f"\nSaved {OUT.relative_to(REPO)} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
